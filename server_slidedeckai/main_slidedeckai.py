from fastapi import FastAPI, Form, HTTPException, UploadFile, File
from fastapi.responses import FileResponse, JSONResponse
from contextlib import asynccontextmanager
from PIL import Image
import os
import torch
from transformers import pipeline
import gc
from io import BytesIO
from tqdm import tqdm

from json_repair import repair_json

# make sure there is a progress bar
from diffusers.utils.logging import enable_progress_bar
enable_progress_bar()

##################
## SLIDES PART ###
##################

# slow HF models loading during application startup
@asynccontextmanager
async def slidedeckai_lifespan(app: FastAPI):
    app.state.max_tokens = 300
    app.state.last_input = None
    app.state.messages = None
    
    try:
        print("Loading slide deck model (gemma-3-12b-it)...")
        slide_model_id = "gemma-3-12b-it"
        # for some reason, the GGUF quantized model only supports text generation
        # I decided to just use HF since the above uses HF also, but there seems to be a general recommendation for Ollama
        # feel free to change to Ollama if already using Ollama to run a LLM just need to change prompt
        slides_pipe = pipeline(
            "text-generation",
            model=slide_model_id,
            device="cuda",
            torch_dtype=torch.bfloat16
        )
        app.state.slides_pipe = slides_pipe
        app.state.slides_template = None
        ### ADD SYSTEM PROMPT FOR JSON GENERATION FOR PPTX ### 
        app.state.messages = [{
            "role": "system",
            "content": [{"type": "text", "text": system_prompt}]
        }]
        print("Slides model (gemma-3-12b-it) loaded successfully.")
    except Exception as e:
        app.state.slides_pipe = None
        app.state.messages = None
        print(f"Slides model (gemma-3-12b-it) loading failed: {e}")

    yield  # App is running

    # Cleanup here
    try:
        if hasattr(app.state, "slides_pipe"):
            del app.state.slides_pipe
    except Exception as e:
        print(f"Error during cleanup: {e}")

    gc.collect()

    torch.cuda.empty_cache()
    torch.cuda.ipc_collect()
    print("Memory fully cleared.")

# app = FastAPI(lifespan=lifespan)
app = FastAPI(lifespan=slidedeckai_lifespan)

########################
## NORMAL DEBUG STUFF ##
########################
# Health Check 
@app.get("/")
def root():
    return {"message": "Slide Gen ok"}

@app.get("/debug_on")
def debug_on():
    try:
        app.state.max_tokens = 20
        return {"message": "Slide Gen Debug ON"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error setting Slide Gen debug mode ON")

@app.get("/debug_off")
def debug_off():
    try:
        app.state.max_tokens = 300
        return {"message": "Slide Gen Debug OFF"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error setting Slide Gen debug mode OFF")
    
########################
## SLIDE DECK AI PART ##
########################

# slides portion of API (receive user_input)
@app.post("/input_slides")
async def receive_input_slides(user_input: str = Form(...)):
    ### ADD USER MESSAGE TO THE OBJECT ###
    try:
        user_message = {
            "role": "user",
            ### to add the text from PDF into the content
            "content": [{"type": "text", "text": user_input}]
        }
        app.state.messages.append(user_message)
        print(f"Received slides prompt: {user_input}")
        return {"message": f"Slide Gen Prompt received: {user_input}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error setting slides generation prompt")

@app.post("/input_template")
async def receive_input_template(user_slides: str = Form(...)):
    ### ADD USER MESSAGE TO THE OBJECT ###
    try:
        user_slides = app.state.slides_template
        return {"message": f"Slide Template received: {user_slides}"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error receiving slides template name")

# slides portion of API (generate compliant JSON from LLM)
@app.get("/slides_json")
def get_slides_json():
    if not app.state.last_input:
        raise HTTPException(status_code=400, detail="No prompt provided yet. Call /input_slides first.")
    if not app.state.slides_pipe:
        raise HTTPException(status_code=500, detail="Slides model not loaded.")
    ### GET INFERENCE FROM GEMMA USING THE MESSAGES OBJECT ###
    raw_json_output = app.state.slides_pipe(text=app.state.messages, max_new_tokens=300)
    json_output = repair_json(raw_json_output)
    return JSONResponse(content=json_output)

#TODO implement the template thing
# slides portion of API (convert JSON to slides and get it)
@app.get("/slides_from_json")
async def slides_from_json(json_data: str = Form(...), template: str = Form(...)):
    try:
        import json
        # Parse the JSON string
        slides_data = json.loads(json_data)


        if template in PPTX_TEMPLATE_FILES:
            template_path = PPTX_TEMPLATE_FILES[template]['file']
            if not os.path.isfile(template):
                print(f"Template file '{template_path}' for '{template}' not found. Using default template.")
                template = None
        # Name not valid
        else:
            available_names = ', '.join(PPTX_TEMPLATE_FILES.keys())
            print(f"Template '{template}' not recognized as name or valid file path. Available names: {available_names}. Using default template.")
            template = None
        
        os.makedirs("generated_pptx", exist_ok=True)

        # Create unique filename (to avoid overwriting)
        filename = f"presentation_test.pptx"
        save_path = os.path.join("generated_pptx", filename)
        
        ### in generate powerpoint presentation have a argument to accept the name of the template used and then resolve to a file path
        headers = generate_powerpoint_presentation(slides_data, save_path)
        print(f"Generated Slide Headers: {headers}")
       
        # TODO whats the media type for pptx
        return FileResponse(save_path, media_type="image/png")
        
    except json.JSONDecodeError:
        raise HTTPException(status_code=400, detail="Invalid JSON data")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")

### READ PDF TO ADD TO MESSAGES OBJECT PART ###
### SUGGEST LOOKING INTO PyPdf library like in Slide Deck AI to convert PDF to text ###
#TODO implement the reading the PDF for additional info thing
# @app.post("/input_pdf")
# def input_pdf(pdf: UploadFile = File(...)):
#     return

### Use this function as reference for input_pdf
# def get_pdf_contents(
#         pdf_file: st.runtime.uploaded_file_manager.UploadedFile,
#         page_range: tuple[int, int]) -> str:
#     """
#     Extract the text contents from a PDF file.

#     :param pdf_file: The uploaded PDF file.
#     :param page_range: The range of pages to extract contents from.
#     :return: The contents.
#     """

#     reader = PdfReader(pdf_file)

#     start, end = page_range  # Set start and end per the range (user-specified values)

#     text = ''

#     if end is None:
#         # If end is None (where PDF has only 1 page or start = end), extract start
#         end = start

#     # Get the text from the specified page range
#     for page_num in range(start - 1, end):
#         text += reader.pages[page_num].extract_text()

#     return text
#####################
## PROMPT CONSTANT ##
#####################

### add a instruction to the LLM to accept text from the PDF document and the topic from the user
system_prompt = """
You are a helpful, intelligent assistant. You are experienced with PowerPoint.

Create the slides for a presentation on the given topic and text from a PDF document.

<the text is in <text></text>
Include main headings for each slide, detailed bullet points for each slide.
Add relevant, detailed content to each slide. When relevant, add one or two EXAMPLES to illustrate the concept.
For two or three important slides, generate the key message that those slides convey.
Present numbers/facts in slides with tables whenever applicable.
Any slide with a table must not have any other content such as bullet points.
E.g., you can tabulate data to summarize some facts on the topic, metrics, experimental settings/results, compare features, and so on.

Read this information carefully. Based on the contents provided, organize the presentation.
For example, if it's a paper, you can consider having slides describing Problem, Solution, Experiments, and Results, among other sections.
If it's a product brochure, you can have Features, Changes, Operating Conditions, and likewise relevant sections.
Similarly, decide for other content types. Then appropriately incorporate the contents into the relevant slides, presenting in a useful way.
If there are important content, e.g., equations and theorems, try to capture a few of them.
Overall, rather than creating a bulleted list of all information, present them in a meaningful way.

Identify if a slide describes a step-by-step/sequential process, then begin the bullet points with a special marker >>.
Limit this to max two or three slides.

Also, add at least one slide with a double column layout by generating appropriate content based on the description in the JSON schema provided below.
In addition, for each slide, add image keywords based on the content of the respective slides.
These keywords will be later used to search for images from the Web relevant to the slide content.

In addition, create one slide containing 4 TO 6 icons (pictograms) illustrating some key ideas/aspects/concepts relevant to the topic.
In this slide, each line of text will begin with the name of a relevant icon enclosed between [[ and ]], e.g., [[machine-learning]] and [[fairness]].
Insert icons only in this slide. Icon names must not be Unicode emojis.

The content of each slide should be detailed and descriptive but not way too verbose.
Avoid writing like a report, but also avoid very short bullet points with just 3-4 words.
Each bullet point should be detailed and explanatory, not just short phrases.
You can use Markdown-like styles for bold & italics.

ALWAYS add a concluding slide at the end, containing a list of the key takeaways and an optional call-to-action if relevant to the context.
Unless explicitly instructed with the topic, create 10 to 12 slides. You must never create more than 15 to 20 slides.

When possible, try to create the slides in the same language as the topic.
`img_keywords` MUST always be in English.

The output must be only a valid and syntactically correct JSON adhering to the following schema:
{
    "title": "Presentation Title",
    "slides": [
        {
            "heading": "Heading for the First Slide",
            "bullet_points": [
                "First bullet point",
                [
                    "Sub-bullet point 1",
                    "Sub-bullet point 2"
                ],
                "Second bullet point"
            ],
            "key_message": "",
            "img_keywords": "a few keywords"
        },
        {
            "heading": "Heading for the Second Slide",
            "bullet_points": [
                "First bullet point",
                "Second bullet item",
                "Third bullet point"
            ],
            "key_message": "The key message conveyed in this slide",
            "img_keywords": "some keywords for this slide"
        },
        {
            "heading": "A slide illustrating key ideas/aspects/concepts (Hint: generate an appropriate heading)",
            "bullet_points": [
                "Some text",
                "Some words describing this aspect",
                "Another aspect highlighted here",
                "Another point here"
            ],
            "key_message": "",
            "img_keywords": ""
        },
        {
            "heading": "A slide that describes a step-by-step/sequential process",
            "bullet_points": [
                ">> The first step of the process (begins with special marker >>)",
                ">> A second step (begins with >>)",
                ">> Third step"
            ],
            "key_message": "",
            "img_keywords": ""
        },
        {
            "heading": "A slide with a double column layout (useful for side-by-side comparison/contrasting of two related concepts, e.g., pros & cons, advantages & risks, old approach vs. modern approach, and so on)",
            "bullet_points": [
                {
                    "heading": "Heading of the left column",
                    "bullet_points": [
                        "First bullet point",
                        "Second bullet item",
                        "Third bullet point"
                    ]
                },
                {
                    "heading": "Heading of the right column",
                    "bullet_points": [
                        "First bullet point",
                        "Second bullet item",
                        "Third bullet point"
                    ]
                }
            ],
            "key_message": "",
            "img_keywords": ""
        },
        {
            "heading": "Slide with a table",
            "table": {
                "headers": ["Column 1", "Column 2", "Column 3"],
                "rows": [
                    ["Row 1, Col 1", "Row 1, Col 2", "Row 1, Col 3"],
                    ["Row 2, Col 1", "Row 2, Col 2", "Row 2, Col 3"],
                    ["Row 3, Col 1", "Row 3, Col 2", "Row 3, Col 3"]
                ]
            },
            "key_message": "",
            "img_keywords": "leave empty"
        }
    ]
}
"""

#############################
## TEMPLATE FILE CONSTANTS ##
##############################

PPTX_TEMPLATE_FILES = {
    'Basic': {
        'file': 'pptx_templates/Blank.pptx',
    },
    'Ion Boardroom': {
        'file': 'pptx_templates/Ion_Boardroom.pptx',
    },
    'Minimalist Sales Pitch': {
        'file': 'pptx_templates/Minimalist_sales_pitch.pptx',
    },
    'Urban Monochrome': {
        'file': 'pptx_templates/Urban_monochrome.pptx',
    },
}


#################################################################
## CLASS, MEHTODS AND FUNCTIONS TO CONVERT JSON TO PPTX (TEST) ##
#################################################################
"""
A set of functions to create a PowerPoint slide deck - CLEANED VERSION
"""
# this entire class, methods and functions are reverse engineered from Slide Deck AI and cleaned by Claude Sonnet 4
# https://github.com/barun-saha/slide-deck-ai/blob/main/helpers/pptx_helper.py
# if it doesn't work idk

import json
import logging
import os
import re
from typing import List, Tuple, Dict, Any

import pptx
# from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

# Constants
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

# Markers and patterns
STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'
SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')

# Display settings
ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1
IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

# Colors
ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]

logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)


class PresentationData:
    """Clean data structure for presentation content."""
    
    def __init__(self, data: Dict[str, Any]):
        self.title = data.get('title', 'Untitled Presentation')
        self.slides = [SlideData(slide) for slide in data.get('slides', [])]


class SlideData:
    """Clean data structure for slide content."""
    
    def __init__(self, data: Dict[str, Any]):
        self.heading = data.get('heading', '')
        self.bullet_points = data.get('bullet_points', [])
        self.key_message = data.get('key_message', '')
        self.img_keywords = data.get('img_keywords', '')
        self.table = data.get('table')
    
    @property
    def clean_heading(self) -> str:
        """Remove slide numbers from heading."""
        if SLIDE_NUMBER_REGEX.match(self.heading):
            idx = self.heading.find(':')
            return self.heading[idx + 1:].strip()
        return self.heading


def load_presentation_data(data_source) -> PresentationData:
    """
    Load presentation data from various sources with proper validation.
    
    :param data_source: Can be a dict, JSON string, or file path
    :return: PresentationData object
    """
    if isinstance(data_source, dict):
        return PresentationData(data_source)
    
    if isinstance(data_source, str):
        try:
            # Try to parse as JSON string first
            data = json.loads(data_source)
            return PresentationData(data)
        except json.JSONDecodeError:
            # Assume it's a file path
            if os.path.exists(data_source):
                with open(data_source, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return PresentationData(data)
            else:
                raise FileNotFoundError(f"File not found: {data_source}")
    
    raise ValueError("Invalid data source type")


def format_text_in_paragraph(paragraph, text: str):
    """Apply bold and italic formatting to text in a paragraph."""
    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0
    
    for match in matches:
        start, end = match.span()
        
        # Add unformatted text before the formatted section
        if start > last_index:
            run = paragraph.add_run()
            run.text = text[last_index:start]
        
        # Extract formatted text
        if match.group(2):  # Bold
            run = paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Italics
            run = paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True
        
        last_index = end
    
    # Add any remaining unformatted text
    if last_index < len(text):
        run = paragraph.add_run()
        run.text = text[last_index:]


def get_flat_bullet_points(items: list, level: int = 0) -> List[Tuple[str, int]]:
    """
    Flatten hierarchical bullet points into a single list with levels.
    
    :param items: List of bullet points (strings or nested lists)
    :param level: Current hierarchy level
    :return: List of (text, level) tuples
    """
    flat_list = []
    
    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list.extend(get_flat_bullet_points(item, level + 1))
        elif isinstance(item, dict):
            # Handle special cases like double column layouts
            flat_list.append((item, level))
    
    return flat_list


def add_bullet_points_to_frame(text_frame, bullet_points: List[Tuple[str, int]]):
    """Add bullet points to a text frame with proper formatting."""
    for idx, (text, level) in enumerate(bullet_points):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.level = level
        
        clean_text = text.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        format_text_in_paragraph(paragraph, clean_text)


class PowerPointGenerator:
    """Main class for generating PowerPoint presentations."""
    
    def __init__(self, template_path: str = None):
        ### where template is added 
        self.template_path = template_path or self._get_default_template()
    
    def _get_default_template(self) -> str:
        """Get default template path - override this method as needed."""
        return None  # Uses default PowerPoint template
    
    def create_presentation(self, presentation_data: PresentationData, output_path: str) -> List[str]:
        """
        Create a PowerPoint presentation from structured data.
        
        :param presentation_data: The presentation content
        :param slides_template: The PPTX template to use.
        :param output_path: Output file path
        :return: List of slide headers
        """
        if self.template_path:
            presentation = pptx.Presentation(self.template_path)
        else:
            presentation = pptx.Presentation()
        
        slide_width_inch, slide_height_inch = self._get_slide_dimensions(presentation)
        headers = []
        
        # Title slide
        headers.append(self._create_title_slide(presentation, presentation_data.title))
        
        # Content slides
        for slide_data in presentation_data.slides:
            try:
                self._create_content_slide(presentation, slide_data, slide_width_inch, slide_height_inch)
                headers.append(slide_data.clean_heading)
            except Exception as e:
                logger.error(f"Error creating slide '{slide_data.heading}': {e}")
                continue
        
        # Thank you slide
        self._create_thank_you_slide(presentation)
        
        # Save presentation
        presentation.save(output_path)
        logger.info(f"Presentation saved: {output_path}")
        
        return headers
    
    def _get_slide_dimensions(self, presentation) -> Tuple[float, float]:
        """Get slide dimensions in inches."""
        slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
        slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
        return slide_width_inch, slide_height_inch
    
    def _create_title_slide(self, presentation, title: str) -> str:
        """Create the title slide."""
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = 'Generated with Python-PPTX'
        
        return title
    
    def _create_thank_you_slide(self, presentation):
        """Create a thank you slide."""
        last_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(last_slide_layout)
        slide.shapes.title.text = 'Thank you!'
    
    def _create_content_slide(self, presentation, slide_data: SlideData, width: float, height: float):
        """Create a content slide with appropriate layout."""
        # Try different slide types based on content
        if self._try_create_table_slide(presentation, slide_data):
            return
        
        if self._try_create_icon_slide(presentation, slide_data, width, height):
            return

        ### Not handled yet
        if self._try_create_process_slide(presentation, slide_data, width, height):
            return
        
        if self._try_create_comparison_slide(presentation, slide_data, width, height):
            return
        
        # Default bullet point slide
        self._create_default_slide(presentation, slide_data, width, height)
    
    def _try_create_table_slide(self, presentation, slide_data: SlideData) -> bool:
        """Try to create a table slide if table data exists."""
        if not slide_data.table:
            return False
        # if slide_data.image_keywords:
        #     search(slide_data.image_keywords)
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        headers = slide_data.table.get('headers', [])
        rows = slide_data.table.get('rows', [])
        
        if not headers or not rows:
            return False
        
        # Position table
        placeholder = slide.placeholders[1]
        table = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            placeholder.left, placeholder.top,
            placeholder.width, placeholder.height
        ).table
        
        # Set headers
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Fill rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(headers):  # Prevent index errors
                    table.cell(row_idx, col_idx).text = str(cell_text)
        
        return True
    
    def _try_create_icon_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create an icon-based slide."""
        if not slide_data.bullet_points:
            return False
        # if slide_data.image_keywords:
        #     search(slide_data.image_keywords)
        # Check if all bullet points are icon markers
        for item in slide_data.bullet_points:
            if not isinstance(item, str) or not item.startswith(ICON_BEGINNING_MARKER):
                return False
        
        # This would need icon handling logic
        # For now, fall back to default
        return False
    
    def _try_create_process_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create a step-by-step process slide."""
        if not slide_data.bullet_points:
            return False
        # if slide_data.image_keywords:
        #     search(slide_data.image_keywords)        
        steps = slide_data.bullet_points
        step_count = 0
        
        # Count steps with process markers
        for step in steps:
            if isinstance(step, str) and step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                step_count += 1
        
        # Need at least 3 steps and most should have markers
        if step_count < 3 or step_count / len(steps) < 0.5:
            return False
        
        # Create process slide layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        # This would need process visualization logic
        # For now, create as bullet points
        placeholder = slide.placeholders[1]
        flat_points = get_flat_bullet_points(steps)
        add_bullet_points_to_frame(placeholder.text_frame, flat_points)
        
        return True
    
    def _try_create_comparison_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create a comparison slide with two columns."""
        if not slide_data.bullet_points or len(slide_data.bullet_points) != 2:
            return False
        # if slide_data.image_keywords:
        #     search(slide_data.image_keywords)   
        # Check if we have two dictionary objects for comparison
        left_col, right_col = slide_data.bullet_points[0], slide_data.bullet_points[1]
        if not (isinstance(left_col, dict) and isinstance(right_col, dict)):
            return False
        
        try:
            slide = presentation.slides.add_slide(presentation.slide_layouts[4])  # Two content layout
        except IndexError:
            return False  # Layout not available
        
        slide.shapes.title.text = slide_data.clean_heading
        
        # Get placeholders - try standard indices first, then search by name
        left_heading, right_heading, left_content, right_content = self._get_comparison_placeholders(slide)
        
        if not (left_content and right_content):
            logger.warning("Could not find required placeholders for comparison slide")
            return False
        
        # Handle left column
        if 'heading' in left_col and left_heading:
            left_heading.text = left_col['heading']
        
        if 'bullet_points' in left_col:
            flat_items = get_flat_bullet_points(left_col['bullet_points'], level=0)
            
            # If no separate heading placeholder, add heading to content
            if not left_heading and 'heading' in left_col:
                left_content.text_frame.text = left_col['heading']
            
            add_bullet_points_to_frame(left_content.text_frame, flat_items)
        
        # Handle right column
        if 'heading' in right_col and right_heading:
            right_heading.text = right_col['heading']
        
        if 'bullet_points' in right_col:
            flat_items = get_flat_bullet_points(right_col['bullet_points'], level=0)
            
            # If no separate heading placeholder, add heading to content
            if not right_heading and 'heading' in right_col:
                right_content.text_frame.text = right_col['heading']
            
            add_bullet_points_to_frame(right_content.text_frame, flat_items)
        
        # Add key message if present
        if slide_data.key_message:
            self._add_key_message(slide, slide_data.key_message, width, height)
        
        return True
    
    def _create_default_slide(self, presentation, slide_data: SlideData, width: float, height: float):
        """Create a default bullet point slide."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        try:
            body_placeholder = slide.placeholders[1]
        except (IndexError, KeyError):
            # Find first content placeholder
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 2:  # Content placeholder
                    body_placeholder = placeholder
                    break
            else:
                logger.warning("No suitable placeholder found for slide content")
                return
        
        if slide_data.bullet_points:
            flat_points = get_flat_bullet_points(slide_data.bullet_points)
            add_bullet_points_to_frame(body_placeholder.text_frame, flat_points)
    
    def _get_comparison_placeholders(self, slide):
        """Get placeholders for two-column comparison layout."""
        shapes = slide.shapes
        left_heading = right_heading = left_content = right_content = None
        
        # Try standard placeholder indices first
        try:
            left_heading = shapes.placeholders[1]
            left_content = shapes.placeholders[2]
            right_heading = shapes.placeholders[3]
            right_content = shapes.placeholders[4]
        except KeyError:
            # If standard indices fail, search by placeholder names
            placeholders = self._get_slide_placeholders(slide, 4)
            
            for idx, name in placeholders:
                name_lower = name.lower()
                
                if 'text placeholder' in name_lower:
                    if not left_heading:
                        left_heading = shapes.placeholders[idx]
                    elif not right_heading:
                        right_heading = shapes.placeholders[idx]
                
                elif 'content placeholder' in name_lower:
                    if not left_content:
                        left_content = shapes.placeholders[idx]
                    elif not right_content:
                        right_content = shapes.placeholders[idx]
        
        return left_heading, right_heading, left_content, right_content
    
    def _get_slide_placeholders(self, slide, layout_number: int):
        """Get placeholder info for a slide, excluding title placeholder."""
        placeholders = [
            (shape.placeholder_format.idx, shape.name.lower()) 
            for shape in slide.shapes.placeholders
            if shape.placeholder_format.idx != 0  # Exclude title
        ]
        return placeholders
    
    def _add_key_message(self, slide, key_message: str, slide_width_inch: float, slide_height_inch: float):
        """Add a key message shape to the slide."""
        if not key_message.strip():
            return
        
        height = pptx.util.Inches(1.6)
        width = pptx.util.Inches(slide_width_inch / 2.3)
        top = pptx.util.Inches(slide_height_inch - height.inches - 0.1)
        left = pptx.util.Inches((slide_width_inch - width.inches) / 2)
        
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        
        format_text_in_paragraph(shape.text_frame.paragraphs[0], key_message)

def generate_powerpoint_presentation(
    data_source,
    output_file_path: str = None,
    template_path: str = None
) -> List[str]:
    """
    Main function to generate PowerPoint presentation.
    
    :param data_source: Presentation data (dict, JSON string, or file path)
    :param output_file_path: Output file path
    :param template_path: Optional template file path
    :return: List of slide headers
    """
    try:
        # Load and validate data
        presentation_data = load_presentation_data(data_source)
        
        # Create generator and build presentation
        generator = PowerPointGenerator(template_path)
        headers = generator.create_presentation(presentation_data, output_file_path)
        
        return headers
        
    except Exception as e:
        logger.error(f"Error generating presentation: {e}")
        raise
