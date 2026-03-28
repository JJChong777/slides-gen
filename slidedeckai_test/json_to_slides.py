"""
A set of functions to create a PowerPoint slide deck - CLEANED VERSION
"""
# this entire class, methods and functions are reverse engineered from Slide Deck AI and cleaned by Claude Sonnet 4
# https://github.com/barun-saha/slide-deck-ai/blob/main/helpers/pptx_helper.py
# if it doesn't work idk

import json
import logging
import os
# import pathlib
import re
from typing import List, Tuple, Dict, Any

import pptx
# from dotenv import load_dotenv
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.placeholder import PicturePlaceholder, SlidePlaceholder

# Constants
EMU_TO_INCH_SCALING_FACTOR = 1.0 / 914400
INCHES_3 = pptx.util.Inches(3)
INCHES_2 = pptx.util.Inches(2)
INCHES_1_5 = pptx.util.Inches(1.5)
INCHES_1 = pptx.util.Inches(1)
INCHES_0_8 = pptx.util.Inches(0.8)
INCHES_0_9 = pptx.util.Inches(0.9)
INCHES_0_5 = pptx.util.Inches(0.5)
INCHES_0_4 = pptx.util.Inches(0.4)
INCHES_0_3 = pptx.util.Inches(0.3)
INCHES_0_2 = pptx.util.Inches(0.2)

# Markers and patterns
STEP_BY_STEP_PROCESS_MARKER = '>> '
ICON_BEGINNING_MARKER = '[['
ICON_END_MARKER = ']]'
SLIDE_NUMBER_REGEX = re.compile(r"^slide[ ]+\d+:", re.IGNORECASE)
ICONS_REGEX = re.compile(r"\[\[(.*?)\]\]\s*(.*)")
BOLD_ITALICS_PATTERN = re.compile(r'(\*\*(.*?)\*\*|\*(.*?)\*)')

# Display settings
ICON_SIZE = INCHES_0_8
ICON_BG_SIZE = INCHES_1
IMAGE_DISPLAY_PROBABILITY = 1 / 3.0
FOREGROUND_IMAGE_PROBABILITY = 0.8

# Colors
ICON_COLORS = [
    pptx.dml.color.RGBColor.from_string('800000'),  # Maroon
    pptx.dml.color.RGBColor.from_string('6A5ACD'),  # SlateBlue
    pptx.dml.color.RGBColor.from_string('556B2F'),  # DarkOliveGreen
    pptx.dml.color.RGBColor.from_string('2F4F4F'),  # DarkSlateGray
    pptx.dml.color.RGBColor.from_string('4682B4'),  # SteelBlue
    pptx.dml.color.RGBColor.from_string('5F9EA0'),  # CadetBlue
]

logger = logging.getLogger(__name__)
logging.getLogger('PIL.PngImagePlugin').setLevel(logging.ERROR)


class PresentationData:
    """Clean data structure for presentation content."""
    
    def __init__(self, data: Dict[str, Any]):
        self.title = data.get('title', 'Untitled Presentation')
        self.slides = [SlideData(slide) for slide in data.get('slides', [])]


class SlideData:
    """Clean data structure for slide content."""
    
    def __init__(self, data: Dict[str, Any]):
        self.heading = data.get('heading', '')
        self.bullet_points = data.get('bullet_points', [])
        self.key_message = data.get('key_message', '')
        self.img_keywords = data.get('img_keywords', '')
        self.table = data.get('table')
    
    @property
    def clean_heading(self) -> str:
        """Remove slide numbers from heading."""
        if SLIDE_NUMBER_REGEX.match(self.heading):
            idx = self.heading.find(':')
            return self.heading[idx + 1:].strip()
        return self.heading


def load_presentation_data(data_source) -> PresentationData:
    """
    Load presentation data from various sources with proper validation.
    
    :param data_source: Can be a dict, JSON string, or file path
    :return: PresentationData object
    """
    if isinstance(data_source, dict):
        return PresentationData(data_source)
    
    if isinstance(data_source, str):
        try:
            # Try to parse as JSON string first
            data = json.loads(data_source)
            return PresentationData(data)
        except json.JSONDecodeError:
            # Assume it's a file path
            if os.path.exists(data_source):
                with open(data_source, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                return PresentationData(data)
            else:
                raise FileNotFoundError(f"File not found: {data_source}")
    
    raise ValueError("Invalid data source type")


def format_text_in_paragraph(paragraph, text: str):
    """Apply bold and italic formatting to text in a paragraph."""
    matches = list(BOLD_ITALICS_PATTERN.finditer(text))
    last_index = 0
    
    for match in matches:
        start, end = match.span()
        
        # Add unformatted text before the formatted section
        if start > last_index:
            run = paragraph.add_run()
            run.text = text[last_index:start]
        
        # Extract formatted text
        if match.group(2):  # Bold
            run = paragraph.add_run()
            run.text = match.group(2)
            run.font.bold = True
        elif match.group(3):  # Italics
            run = paragraph.add_run()
            run.text = match.group(3)
            run.font.italic = True
        
        last_index = end
    
    # Add any remaining unformatted text
    if last_index < len(text):
        run = paragraph.add_run()
        run.text = text[last_index:]


def get_flat_bullet_points(items: list, level: int = 0) -> List[Tuple[str, int]]:
    """
    Flatten hierarchical bullet points into a single list with levels.
    
    :param items: List of bullet points (strings or nested lists)
    :param level: Current hierarchy level
    :return: List of (text, level) tuples
    """
    flat_list = []
    
    for item in items:
        if isinstance(item, str):
            flat_list.append((item, level))
        elif isinstance(item, list):
            flat_list.extend(get_flat_bullet_points(item, level + 1))
        elif isinstance(item, dict):
            # Handle special cases like double column layouts
            flat_list.append((item, level))
    
    return flat_list


def add_bullet_points_to_frame(text_frame, bullet_points: List[Tuple[str, int]]):
    """Add bullet points to a text frame with proper formatting."""
    for idx, (text, level) in enumerate(bullet_points):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
            paragraph.level = level
        
        clean_text = text.removeprefix(STEP_BY_STEP_PROCESS_MARKER)
        format_text_in_paragraph(paragraph, clean_text)


class PowerPointGenerator:
    """Main class for generating PowerPoint presentations."""
    
    def __init__(self, template_path: str = None):
        self.template_path = template_path or self._get_default_template()
    
    def _get_default_template(self) -> str:
        """Get default template path - override this method as needed."""
        return None  # Uses default PowerPoint template
    
    def create_presentation(self, presentation_data: PresentationData, output_path: str) -> List[str]:
        """
        Create a PowerPoint presentation from structured data.
        
        :param presentation_data: The presentation content
        :param slides_template: The PPTX template to use.
        :param output_path: Output file path
        :return: List of slide headers
        """
        if self.template_path:
            presentation = pptx.Presentation(self.template_path)
        else:
            presentation = pptx.Presentation()
        
        slide_width_inch, slide_height_inch = self._get_slide_dimensions(presentation)
        headers = []
        
        # Title slide
        headers.append(self._create_title_slide(presentation, presentation_data.title))
        
        # Content slides
        for slide_data in presentation_data.slides:
            try:
                self._create_content_slide(presentation, slide_data, slide_width_inch, slide_height_inch)
                headers.append(slide_data.clean_heading)
            except Exception as e:
                logger.error(f"Error creating slide '{slide_data.heading}': {e}")
                continue
        
        # Thank you slide
        self._create_thank_you_slide(presentation)
        
        # Save presentation
        presentation.save(output_path)
        logger.info(f"Presentation saved: {output_path}")
        
        return headers
    
    def _get_slide_dimensions(self, presentation) -> Tuple[float, float]:
        """Get slide dimensions in inches."""
        slide_width_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_width
        slide_height_inch = EMU_TO_INCH_SCALING_FACTOR * presentation.slide_height
        return slide_width_inch, slide_height_inch
    
    def _create_title_slide(self, presentation, title: str) -> str:
        """Create the title slide."""
        title_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(title_slide_layout)
        
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = 'Generated with Python-PPTX'
        
        return title
    
    def _create_thank_you_slide(self, presentation):
        """Create a thank you slide."""
        last_slide_layout = presentation.slide_layouts[0]
        slide = presentation.slides.add_slide(last_slide_layout)
        slide.shapes.title.text = 'Thank you!'
    
    def _create_content_slide(self, presentation, slide_data: SlideData, width: float, height: float):
        """Create a content slide with appropriate layout."""
        # Try different slide types based on content
        if self._try_create_table_slide(presentation, slide_data):
            return
        
        if self._try_create_icon_slide(presentation, slide_data, width, height):
            return
        
        if self._try_create_process_slide(presentation, slide_data, width, height):
            return
        
        if self._try_create_comparison_slide(presentation, slide_data, width, height):
            return
        
        # Default bullet point slide
        self._create_default_slide(presentation, slide_data, width, height)
    
    def _try_create_table_slide(self, presentation, slide_data: SlideData) -> bool:
        """Try to create a table slide if table data exists."""
        if not slide_data.table:
            return False
        
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        headers = slide_data.table.get('headers', [])
        rows = slide_data.table.get('rows', [])
        
        if not headers or not rows:
            return False
        
        # Position table
        placeholder = slide.placeholders[1]
        table = slide.shapes.add_table(
            len(rows) + 1, len(headers),
            placeholder.left, placeholder.top,
            placeholder.width, placeholder.height
        ).table
        
        # Set headers
        for col_idx, header_text in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header_text
            cell.text_frame.paragraphs[0].font.bold = True
        
        # Fill rows
        for row_idx, row_data in enumerate(rows, start=1):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < len(headers):  # Prevent index errors
                    table.cell(row_idx, col_idx).text = str(cell_text)
        
        return True
    
    def _try_create_icon_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create an icon-based slide."""
        if not slide_data.bullet_points:
            return False
        
        # Check if all bullet points are icon markers
        for item in slide_data.bullet_points:
            if not isinstance(item, str) or not item.startswith(ICON_BEGINNING_MARKER):
                return False
        
        # This would need icon handling logic
        # For now, fall back to default
        return False
    
    def _try_create_process_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create a step-by-step process slide."""
        if not slide_data.bullet_points:
            return False
        
        steps = slide_data.bullet_points
        step_count = 0
        
        # Count steps with process markers
        for step in steps:
            if isinstance(step, str) and step.startswith(STEP_BY_STEP_PROCESS_MARKER):
                step_count += 1
        
        # Need at least 3 steps and most should have markers
        if step_count < 3 or step_count / len(steps) < 0.5:
            return False
        
        # Create process slide layout
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        # This would need process visualization logic
        # For now, create as bullet points
        placeholder = slide.placeholders[1]
        flat_points = get_flat_bullet_points(steps)
        add_bullet_points_to_frame(placeholder.text_frame, flat_points)
        
        return True
    
    def _try_create_comparison_slide(self, presentation, slide_data: SlideData, width: float, height: float) -> bool:
        """Try to create a comparison slide with two columns."""
        if not slide_data.bullet_points or len(slide_data.bullet_points) != 2:
            return False
        
        # Check if we have two dictionary objects for comparison
        left_col, right_col = slide_data.bullet_points[0], slide_data.bullet_points[1]
        if not (isinstance(left_col, dict) and isinstance(right_col, dict)):
            return False
        
        try:
            slide = presentation.slides.add_slide(presentation.slide_layouts[4])  # Two content layout
        except IndexError:
            return False  # Layout not available
        
        slide.shapes.title.text = slide_data.clean_heading
        
        # Get placeholders - try standard indices first, then search by name
        left_heading, right_heading, left_content, right_content = self._get_comparison_placeholders(slide)
        
        if not (left_content and right_content):
            logger.warning("Could not find required placeholders for comparison slide")
            return False
        
        # Handle left column
        if 'heading' in left_col and left_heading:
            left_heading.text = left_col['heading']
        
        if 'bullet_points' in left_col:
            flat_items = get_flat_bullet_points(left_col['bullet_points'], level=0)
            
            # If no separate heading placeholder, add heading to content
            if not left_heading and 'heading' in left_col:
                left_content.text_frame.text = left_col['heading']
            
            add_bullet_points_to_frame(left_content.text_frame, flat_items)
        
        # Handle right column
        if 'heading' in right_col and right_heading:
            right_heading.text = right_col['heading']
        
        if 'bullet_points' in right_col:
            flat_items = get_flat_bullet_points(right_col['bullet_points'], level=0)
            
            # If no separate heading placeholder, add heading to content
            if not right_heading and 'heading' in right_col:
                right_content.text_frame.text = right_col['heading']
            
            add_bullet_points_to_frame(right_content.text_frame, flat_items)
        
        # Add key message if present
        if slide_data.key_message:
            self._add_key_message(slide, slide_data.key_message, width, height)
        
        return True
    
    def _create_default_slide(self, presentation, slide_data: SlideData, width: float, height: float):
        """Create a default bullet point slide."""
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = slide_data.clean_heading
        
        try:
            body_placeholder = slide.placeholders[1]
        except (IndexError, KeyError):
            # Find first content placeholder
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.type == 2:  # Content placeholder
                    body_placeholder = placeholder
                    break
            else:
                logger.warning("No suitable placeholder found for slide content")
                return
        
        if slide_data.bullet_points:
            flat_points = get_flat_bullet_points(slide_data.bullet_points)
            add_bullet_points_to_frame(body_placeholder.text_frame, flat_points)
    
    def _get_comparison_placeholders(self, slide):
        """Get placeholders for two-column comparison layout."""
        shapes = slide.shapes
        left_heading = right_heading = left_content = right_content = None
        
        # Try standard placeholder indices first
        try:
            left_heading = shapes.placeholders[1]
            left_content = shapes.placeholders[2]
            right_heading = shapes.placeholders[3]
            right_content = shapes.placeholders[4]
        except KeyError:
            # If standard indices fail, search by placeholder names
            placeholders = self._get_slide_placeholders(slide, 4)
            
            for idx, name in placeholders:
                name_lower = name.lower()
                
                if 'text placeholder' in name_lower:
                    if not left_heading:
                        left_heading = shapes.placeholders[idx]
                    elif not right_heading:
                        right_heading = shapes.placeholders[idx]
                
                elif 'content placeholder' in name_lower:
                    if not left_content:
                        left_content = shapes.placeholders[idx]
                    elif not right_content:
                        right_content = shapes.placeholders[idx]
        
        return left_heading, right_heading, left_content, right_content
    
    def _get_slide_placeholders(self, slide, layout_number: int):
        """Get placeholder info for a slide, excluding title placeholder."""
        placeholders = [
            (shape.placeholder_format.idx, shape.name.lower()) 
            for shape in slide.shapes.placeholders
            if shape.placeholder_format.idx != 0  # Exclude title
        ]
        return placeholders
    
    def _add_key_message(self, slide, key_message: str, slide_width_inch: float, slide_height_inch: float):
        """Add a key message shape to the slide."""
        if not key_message.strip():
            return
        
        height = pptx.util.Inches(1.6)
        width = pptx.util.Inches(slide_width_inch / 2.3)
        top = pptx.util.Inches(slide_height_inch - height.inches - 0.1)
        left = pptx.util.Inches((slide_width_inch - width.inches) / 2)
        
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        
        format_text_in_paragraph(shape.text_frame.paragraphs[0], key_message)

def generate_powerpoint_presentation(
    data_source,
    output_file_path: str = None,
    template_path: str = None
) -> List[str]:
    """
    Main function to generate PowerPoint presentation.
    
    :param data_source: Presentation data (dict, JSON string, or file path)
    :param output_file_path: Output file path
    :param template_path: Optional template file path
    :return: List of slide headers
    """
    try:
        # Load and validate data
        presentation_data = load_presentation_data(data_source)
        
        # Create generator and build presentation
        generator = PowerPointGenerator(template_path)
        headers = generator.create_presentation(presentation_data, output_file_path)
        
        return headers
        
    except Exception as e:
        logger.error(f"Error generating presentation: {e}")
        raise



if __name__ == '__main__':
    # Clean test data - proper JSON without trailing commas
    test_data = {
  "title": "AI Applications: Transforming Industries",
  "slides": [
    {
      "heading": "Introduction to AI Applications",
      "bullet_points": [
        "Artificial Intelligence (AI) is *transforming* various industries",
        "AI applications range from simple decision-making tools to complex systems",
        "AI can be categorized into types: Rule-based, Instance-based, and Model-based"
      ],
      "key_message": "AI is a broad field with diverse applications and categories",
      "img_keywords": "AI, transformation, industries, decision-making, categories"
    },
    {
      "heading": "AI in Everyday Life",
      "bullet_points": [
        "**Virtual assistants** like Siri, Alexa, and Google Assistant",
        "**Recommender systems** in Netflix, Amazon, and Spotify",
        "**Fraud detection** in banking and *credit card* transactions"
      ],
      "key_message": "AI is integrated into our daily lives through various services",
      "img_keywords": "virtual assistants, recommender systems, fraud detection"
    },
    {
      "heading": "AI in Healthcare",
      "bullet_points": [
        "Disease diagnosis and prediction using machine learning algorithms",
        "Personalized medicine and drug discovery",
        "AI-powered robotic surgeries and remote patient monitoring"
      ],
      "key_message": "AI is revolutionizing healthcare with improved diagnostics and patient care",
      "img_keywords": "healthcare, disease diagnosis, personalized medicine, robotic surgeries"
    },
    # this part with the object is really broken
    {
      "heading": "AI in Key Industries",
      "bullet_points": [
        {
          "heading": "Retail",
          "bullet_points": [
            "Inventory management and demand forecasting",
            "Customer segmentation and targeted marketing",
            "AI-driven chatbots for customer service"
          ]
        },
        {
          "heading": "Finance",
          "bullet_points": [
            "Credit scoring and risk assessment",
            "Algorithmic trading and portfolio management",
            "AI for detecting money laundering and cyber fraud"
          ]
        }
      ],
      "key_message": "AI is transforming retail and finance with improved operations and decision-making",
      "img_keywords": "retail, finance, inventory management, credit scoring, algorithmic trading"
    },
    {
      "heading": "AI in Education",
      "bullet_points": [
        "Personalized learning paths and adaptive testing",
        "Intelligent tutoring systems for skill development",
        "AI for predicting student performance and dropout rates"
      ],
      "key_message": "AI is personalizing education and improving student outcomes",
    },
    {
      "heading": "Step-by-Step: AI Development Process",
      "bullet_points": [
        ">> **Step 1:** Define the problem and objectives",
        ">> **Step 2:** Collect and preprocess data",
        ">> **Step 3:** Select and train the AI model",
        ">> **Step 4:** Evaluate and optimize the model",
        ">> **Step 5:** Deploy and monitor the AI system"
      ],
      "key_message": "Developing AI involves a structured process from problem definition to deployment",
      "img_keywords": ""
    },
    # doesn't register the images, which kind of makes sense
    {
      "heading": "AI Icons: Key Aspects",
      "bullet_points": [
        "Human-like *intelligence* and decision-making",
        "Automation and physical *tasks*",
        "Data processing and cloud computing",
        "Insights and *predictions*",
        "Global connectivity and *impact*"
      ],
      "key_message": "AI encompasses various aspects, from human-like intelligence to global impact"
    },
    {
        "heading": "AI vs. ML vs. DL: A Tabular Comparison",
        "table": {
            "headers": ["Feature", "AI", "ML", "DL"],
            "rows": [
                ["Definition", "Creating intelligent agents", "Learning from data", "Deep neural networks"],
                ["Approach", "Rule-based, expert systems", "Algorithms, statistical models", "Deep neural networks"],
                ["Data Requirements", "Varies", "Large datasets", "Massive datasets"],
                ["Complexity", "Varies", "Moderate", "High"],
                ["Computational Cost", "Low to Moderate", "Moderate", "High"],
                ["Examples", "Chess, recommendation systems", "Spam filters, image recognition", "Image recognition, NLP"]
            ]
        },
        "key_message": "This table provides a concise comparison of the key features of AI, ML, and DL.",
        "img_keywords": "AI, ML, DL, comparison, table, features"
    },
    {
      "heading": "Conclusion: Embracing AI's Potential",
      "bullet_points": [
        "AI is transforming industries and improving lives",
        "Ethical considerations are crucial for responsible AI development",
        "Invest in AI education and workforce development",
        "Call to action: Explore AI applications and contribute to shaping its future"
      ],
      "key_message": "AI offers *immense potential*, and we must embrace it responsibly",
      "img_keywords": "AI transformation, ethical considerations, AI education, future of AI"
    },
    {
        "heading": "Tests a step-by-step/sequential process",
        "bullet_points": [
            ">> The first step of the process",
            ">> A second step",
            ">> Third step",
        ],
        "key_message": "",
        "img_keywords": ""
    },
  ]
}
    
    # Generate presentation
    # Make sure the folder exists
    os.makedirs("generated_pptx", exist_ok=True)

    # Create unique filename (to avoid overwriting)
    filename = f"presentation_test.pptx"
    save_path = os.path.join("generated_pptx", filename)

    # Generate presentation
    headers = generate_powerpoint_presentation(test_data, save_path)

    print(f"Generated presentation: {save_path}")
    print(f"Slide headers: {headers}")